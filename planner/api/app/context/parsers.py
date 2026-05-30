"""Format-Parser für Kontext-Dateien (Schritt 2a) — Fließtext + Tabellen.

Reine ephemere Extraktion: bytes -> Text. Kein Persistieren. Unterstützt
.docx, .md, .pdf, .txt, .pptx, .xlsx. Bei nicht unterstütztem Format oder
defektem Inhalt wird `ValueError` geworfen (Router → HTTP 400).
"""

from __future__ import annotations

import io

from ..schemas.project import ContextFormat

_BY_EXT: dict[str, ContextFormat] = {
    "docx": "docx",
    "md": "md",
    "markdown": "md",
    "pdf": "pdf",
    "txt": "txt",
    "text": "txt",
    "pptx": "pptx",
    "xlsx": "xlsx",
}


def detect_format(filename: str) -> ContextFormat:
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    fmt = _BY_EXT.get(ext)
    if fmt is None:
        raise ValueError(f"nicht unterstütztes Format: .{ext or '?'}")
    return fmt


def parse(data: bytes, fmt: ContextFormat) -> str:
    if fmt in ("txt", "md"):
        return data.decode("utf-8", errors="replace")
    if fmt == "docx":
        return _parse_docx(data)
    if fmt == "pptx":
        return _parse_pptx(data)
    if fmt == "xlsx":
        return _parse_xlsx(data)
    if fmt == "pdf":
        return _parse_pdf(data)
    raise ValueError(f"nicht unterstütztes Format: {fmt}")


def _parse_docx(data: bytes) -> str:
    from docx import Document

    doc = Document(io.BytesIO(data))
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _parse_pptx(data: bytes) -> str:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    txt = "".join(run.text for run in para.runs)
                    if txt.strip():
                        parts.append(txt)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    if any(cells):
                        parts.append(" | ".join(cells))
    return "\n".join(parts)


def _parse_xlsx(data: bytes) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    parts: list[str] = []
    try:
        for ws in wb.worksheets:
            parts.append(f"# {ws.title}")
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None]
                if cells:
                    parts.append(" | ".join(cells))
    finally:
        wb.close()
    return "\n".join(parts)


def _parse_pdf(data: bytes) -> str:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(data))
    parts: list[str] = []
    for page in reader.pages:
        txt = page.extract_text() or ""
        if txt.strip():
            parts.append(txt)
    return "\n".join(parts)
