# -*- coding: utf-8 -*-
"""AEGIRA Planner User Guide — konsistentes McK-Deck, from scratch (valide)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

# --- Brand-Tokens (aus dem App-Theme: Inter, Navy/Gold/Steel) ---
NAVY=RGBColor(0x1E,0x27,0x61); NAVYD=RGBColor(0x0E,0x17,0x35)
GOLD=RGBColor(0xF2,0xB3,0x3D); ORANGE=RGBColor(0xE8,0x70,0x3A)
STEEL=RGBColor(0x5B,0x6B,0x85); MUTED=RGBColor(0x8A,0x93,0xA6)
INK=RGBColor(0x14,0x1B,0x34); WHITE=RGBColor(0xFF,0xFF,0xFF)
VELLUM=RGBColor(0xFA,0xF7,0xEE); CARD=RGBColor(0xFF,0xFF,0xFF)
LINE=RGBColor(0xDD,0xE1,0xEA); GREEN=RGBColor(0x2E,0x7D,0x46); RED=RGBColor(0xC3,0x42,0x3F)
F="Inter"; MONO="Consolas"
D="/tmp/deck"

prs=Presentation()
prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
SW,SH=prs.slide_width,prs.slide_height
BLANK=prs.slide_layouts[6]

def slide():
    s=prs.slides.add_slide(BLANK)
    r=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,SW,SH)
    r.fill.solid(); r.fill.fore_color.rgb=VELLUM; r.line.fill.background(); r.shadow.inherit=False
    r._element.addprevious(r._element)  # bg zuerst
    return s

def rect(s,l,t,w,h,fill,line=None,lw=Pt(0.75),rounded=False):
    shp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,l,t,w,h)
    shp.fill.solid(); shp.fill.fore_color.rgb=fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb=line; shp.line.width=lw
    shp.shadow.inherit=False
    return shp

def text(s,l,t,w,h,anchor=MSO_ANCHOR.TOP):
    tb=s.shapes.add_textbox(l,t,w,h); tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    return tf

def line_p(tf,txt,size,color,bold=False,first=False,font=F,after=4,before=0,bullet=None,italic=False):
    p=tf.paragraphs[0] if first else tf.add_paragraph()
    p.space_after=Pt(after); p.space_before=Pt(before)
    if bullet is not None:
        rb=p.add_run(); rb.text=bullet+"  "; rb.font.name=F; rb.font.size=Pt(size); rb.font.bold=True; rb.font.color.rgb=GOLD
    r=p.add_run(); r.text=txt; r.font.name=font; r.font.size=Pt(size); r.font.bold=bold; r.font.italic=italic; r.font.color.rgb=color
    return p

def footer(s,idx,total=24):
    tf=text(s,Inches(0.55),Inches(7.05),Inches(8),Inches(0.35))
    line_p(tf,"exmachinAI · AEGIRA AI Trust Platform · zgpm.aegira.ai",9,MUTED,first=True)
    tf2=text(s,Inches(11.5),Inches(7.05),Inches(1.3),Inches(0.35)); tf2.paragraphs[0].alignment=PP_ALIGN.RIGHT
    line_p(tf2,f"{idx} / {total}",9,MUTED,first=True)

def titlebar(s,kicker,title):
    rect(s,0,0,SW,Inches(1.15),NAVY)
    rect(s,0,Inches(1.15),SW,Inches(0.06),GOLD)
    tf=text(s,Inches(0.6),Inches(0.2),Inches(12.1),Inches(0.85),anchor=MSO_ANCHOR.MIDDLE)
    line_p(tf,kicker,12,GOLD,bold=True,first=True,after=2)
    line_p(tf,title,25,WHITE,bold=True,after=0)

def fit(box,path):
    iw,ih=Image.open(path).size; ar=iw/ih; bl,bt,bw,bh=box
    if bw/bh>ar: h=bh; w=int(bh*ar)
    else: w=bw; h=int(bw/ar)
    return bl+(bw-w)//2, bt+(bh-h)//2, w, h

def screen(s,box,path):
    bl,bt,bw,bh=box
    rect(s,bl-Inches(0.06),bt-Inches(0.06),bw+Inches(0.12),bh+Inches(0.12),WHITE,line=LINE)
    l,t,w,h=fit(box,path)
    pic=s.shapes.add_picture(path,l,t,width=w,height=h)
    pic.line.color.rgb=LINE; pic.line.width=Pt(0.75)

def chip(s,l,t,txt,fill=NAVY,fg=WHITE,w=None):
    w=w or Inches(0.9+0.085*len(txt))
    c=rect(s,l,t,w,Inches(0.32),fill,rounded=True)
    c.text_frame.word_wrap=False; p=c.text_frame.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
    r=p.add_run(); r.text=txt; r.font.name=F; r.font.size=Pt(10.5); r.font.bold=True; r.font.color.rgb=fg
    c.text_frame.vertical_anchor=MSO_ANCHOR.MIDDLE
    return l+w

TOTAL=24
n=[0]
def pg(): n[0]+=1; return n[0]

# ============================================================ 1 TITEL
s=slide()
rect(s,0,0,SW,SH,NAVY)
rect(s,0,Inches(5.0),SW,Inches(0.08),GOLD)
tf=text(s,Inches(0.9),Inches(2.1),Inches(11.5),Inches(3.0))
line_p(tf,"AI TRUST PLATFORM · exmachinAI",14,GOLD,bold=True,first=True,after=14)
line_p(tf,"AEGIRA Agent Operating Model Planner",40,WHITE,bold=True,after=6)
line_p(tf,"User Guide",40,RGBColor(0xC8,0xD2,0xEC),bold=True,after=16)
line_p(tf,"Von der Projektidee zum lauffähigen Agententeam — Schritt für Schritt.",16,RGBColor(0xC8,0xD2,0xEC),after=2)
line_p(tf,"Für Anwender ohne Programmiererfahrung · ZGPM-konform · Stand 05/2026",12,MUTED,after=0)

# ============================================================ 2 ORIENTIERUNG
s=slide(); titlebar(s,"ORIENTIERUNG","Wie du diesen Guide liest")
tf=text(s,Inches(0.6),Inches(1.5),Inches(6.0),Inches(5.2))
line_p(tf,"Jeder der neun Schritte folgt demselben Aufbau — einmal verstanden, findest du dich überall sofort zurecht.",13,STEEL,first=True,after=12)
for lab,desc in [("ZIEL","Was in diesem Schritt erreicht wird."),
                 ("DEINE AKTION","Was du tust — meist lesen, klicken, bestätigen."),
                 ("WAS DAS SYSTEM TUT","Welche Agenten im Hintergrund für dich arbeiten."),
                 ("DEIN KONTROLLPUNKT","Wo du entscheidest. Nichts geht ohne deine Freigabe weiter.")]:
    line_p(tf,lab,12,NAVY,bold=True,after=1,before=4)
    line_p(tf,desc,12.5,INK,after=2)
b=rect(s,Inches(6.9),Inches(1.6),Inches(5.8),Inches(4.9),WHITE,line=LINE,rounded=True)
tf=text(s,Inches(7.2),Inches(1.9),Inches(5.2),Inches(4.4))
line_p(tf,"WAS DU NACH DIESEM GUIDE KANNST",12,ORANGE,bold=True,first=True,after=10)
for i,txt in enumerate(["Ein Projekt in eigenen Worten beschreiben — und methodisch schärfen lassen.",
    "Ein präzises Projektverständnis inkl. Agentenstruktur freigeben.",
    "Einen ZGPM-Plan lesen und anpassen — Meilensteine, RACI, Risiken, Zeit, Kosten.",
    "Die Leitplanken verstehen — was erlaubt ist und was das System verweigert.",
    "Den Agenten-Harness mitgestalten — Artefakte editieren, Struktur per Kommando verbessern.",
    "Ein lauffähiges Paket exportieren — Zip mit Setup für Claude Cowork & Claude Code."],1):
    line_p(tf,txt,12.5,INK,bullet=str(i),after=7)
footer(s,pg())

# ============================================================ 3 DAS PRODUKT
s=slide(); titlebar(s,"DAS PRODUKT","Ein Project-to-Agent-Compiler")
tf=text(s,Inches(0.6),Inches(1.5),Inches(12.1),Inches(0.9))
line_p(tf,"Aus deiner Projektidee entsteht über methodisch saubere Stufen ein vollständiger Bauplan — und daraus ein lauffähiges Agentensystem.",14,STEEL,first=True)
x=Inches(0.6)
for code,title,desc,col in [("PLAN","Projektauftrag → ZGPM-Plan","Meilensteine, Rollen, Risiken.",NAVY),
                            ("HARNESS","Plan → Paket","CLAUDE.md, Skills und Agenten.",ORANGE),
                            ("RUN","Paket läuft","Auf deinem Rechner in Claude Code / Cowork.",GREEN)]:
    c=rect(s,x,Inches(2.7),Inches(3.85),Inches(2.2),WHITE,line=LINE,rounded=True)
    rect(s,x,Inches(2.7),Inches(3.85),Inches(0.12),col)
    tf=text(s,x+Inches(0.3),Inches(3.0),Inches(3.25),Inches(1.7))
    line_p(tf,code,20,col,bold=True,first=True,after=6)
    line_p(tf,title,14,INK,bold=True,after=4)
    line_p(tf,desc,12.5,STEEL,after=0)
    x+=Inches(4.05)
tf=text(s,Inches(0.6),Inches(5.3),Inches(12.1),Inches(1.4))
line_p(tf,"Warum?",13,NAVY,bold=True,first=True,after=3)
line_p(tf,"Wer heute Agenten baut, startet technisch — Prompt, Tool, Workflow. Was fehlt, ist die organisatorische Vorarbeit: Plan, Rollen, RACI, Governance. Solo-Builder haben kein Projektteam. Der Planner liefert genau diese organisatorische Intelligenz.",13,INK,after=0)
footer(s,pg())

# ============================================================ 4 MENTALES MODELL
s=slide(); titlebar(s,"DAS MENTALE MODELL","Agenten sind digitale Mitarbeiter")
tf=text(s,Inches(0.6),Inches(1.5),Inches(12.1),Inches(0.9))
line_p(tf,"Wenn Agenten wie Mitarbeiter arbeiten, brauchen sie, was jedes gute Team braucht: Rollen, Verantwortlichkeiten, Fähigkeiten und Governance.",14,STEEL,first=True)
c=rect(s,Inches(0.6),Inches(2.7),Inches(5.8),Inches(3.4),WHITE,line=LINE,rounded=True)
tf=text(s,Inches(0.95),Inches(3.0),Inches(5.1),Inches(2.8))
line_p(tf,"ÜBLICHER START · TECHNISCH",12,MUTED,bold=True,first=True,after=8)
line_p(tf,"Prompt → Tool → Workflow",15,INK,bold=True,font=MONO,after=10)
line_p(tf,"Ergebnis: technisch lauffähig, aber ohne organisatorische Klarheit.",12.5,STEEL,italic=True,after=0)
c=rect(s,Inches(6.9),Inches(2.7),Inches(5.8),Inches(3.4),NAVY,rounded=True)
rect(s,Inches(6.9),Inches(2.7),Inches(5.8),Inches(0.12),GOLD)
tf=text(s,Inches(7.25),Inches(3.0),Inches(5.1),Inches(2.8))
line_p(tf,"AEGIRA-START · ORGANISATORISCH",12,GOLD,bold=True,first=True,after=8)
line_p(tf,"Projektplan → Rollenmodell → Skill-Matrix → RACI/PVM → Agentenstruktur",15,WHITE,bold=True,after=10)
line_p(tf,"Aus dem Plan wird die Technik abgeleitet — nicht umgekehrt.",12.5,RGBColor(0xC8,0xD2,0xEC),italic=True,after=0)
footer(s,pg())

# ============================================================ 5 DER WEG
s=slide(); titlebar(s,"DER WEG & DEIN PRINZIP","Drei Phasen · neun Schritte · drei Gates")
phases=[("VERSTEHEN","1 Beschreiben · 2 Interview · 3 Verständnis","◆ GATE 1",NAVY),
        ("PLANEN","4 Verwalten · 5 Leitplanken · 6 ZGPM-Plan · 7 Review","◆ GATE 2",ORANGE),
        ("BAUEN","8 Harness · 9 Export","◆ GATE 3",GREEN)]
x=Inches(0.6)
for ph,steps,gate,col in phases:
    c=rect(s,x,Inches(1.6),Inches(3.85),Inches(2.5),WHITE,line=LINE,rounded=True)
    rect(s,x,Inches(1.6),Inches(3.85),Inches(0.5),col)
    tf=text(s,x+Inches(0.25),Inches(1.66),Inches(3.35),Inches(0.4),anchor=MSO_ANCHOR.MIDDLE)
    line_p(tf,ph,14,WHITE,bold=True,first=True)
    tf=text(s,x+Inches(0.25),Inches(2.3),Inches(3.35),Inches(1.3))
    line_p(tf,steps,12.5,INK,first=True,after=0)
    tf=text(s,x+Inches(0.25),Inches(3.55),Inches(3.35),Inches(0.4))
    line_p(tf,gate,12,col,bold=True,first=True)
    x+=Inches(4.05)
c=rect(s,Inches(0.6),Inches(4.45),Inches(12.1),Inches(2.05),NAVY,rounded=True)
tf=text(s,Inches(0.95),Inches(4.7),Inches(11.4),Inches(1.6))
line_p(tf,"DU BEHÄLTST DIE KONTROLLE — HUMAN-IN-THE-LOOP",13,GOLD,bold=True,first=True,after=8)
line_p(tf,"Meilenstein-Freigabe — jeder Phasenübergang wartet auf deine manuelle Bestätigung.",13,WHITE,bullet="✓",after=4)
line_p(tf,"Rote Risiko-Ampel — steht ein Risiko auf Rot, stoppt der Lauf. Nur du gibst ihn frei.",13,WHITE,bullet="!",after=4)
line_p(tf,"Transparenz — du siehst die Agenten arbeiten; jede Aktion landet im Audit-Trail.",13,WHITE,bullet="◆",after=0)
footer(s,pg())

# ============================================================ 6-14 SCHRITTE 1-9
STEPS=[
 ("PHASE VERSTEHEN","SCHRITT 1 VON 9","Beschreibe dein Projekt in eigenen Worten",
  "Das System erfährt, was du vorhast — formlos, in deiner Sprache.",
  "Du schreibst frei drauflos. Keine Fachbegriffe, keine Struktur nötig.",
  "Es liest mit und bereitet die erste Rückfrage vor — nichts wird ohne dich gespeichert.",
  "Du klickst „Weiter“, wenn deine Beschreibung steht.","new"),
 ("PHASE VERSTEHEN","SCHRITT 2 VON 9","Das Schärfungs-Interview (McKinsey-Methode)",
  "Aus der Idee wird ein präzises Verständnis: Projektart, Umfang, Skills.",
  "Du beantwortest gezielte Rückfragen — eine nach der anderen. Optional Dokumente/Dropbox einspeisen.",
  "Es fragt MECE und hypothesengeleitet nach und macht aktiv Vorschläge.",
  "Du kannst jeden Vorschlag annehmen, ändern oder verwerfen.","interview"),
 ("PHASE VERSTEHEN","SCHRITT 3 VON 9","Projektverständnis & Agentenstruktur freigeben",
  "Eine pre-finale Zusammenfassung — plus die Agenten, die deinen Plan bauen.",
  "Du liest, korrigierst Details und gibst frei.",
  "Es verdichtet alles und leitet ab, welche Planungs-Agenten nötig sind.",
  "Ohne deine Freigabe (Gate 1) startet keine Planung.","understanding"),
 ("PHASE PLANEN","SCHRITT 4 VON 9","Deine Projekte verwalten",
  "Jedes Projekt bleibt erhalten — in jeder Phase, jederzeit auffindbar.",
  "Öffnen, umbenennen, als Vorlage duplizieren oder löschen — alles mit einem Klick.",
  "Es speichert jede Phase als eigene Version und zeigt den Status.",
  "Löschen wird immer noch einmal von dir bestätigt.","dashboard"),
 ("PHASE PLANEN","SCHRITT 5 VON 9","Die Leitplanken: was geht — und was nicht",
  "Du siehst transparent, was das System verweigert und warum.",
  "Du liest das Urteil (erlaubt / eskalieren / abgelehnt) und quittierst bewusst.",
  "Es prüft das Vorhaben gegen definierte Kategorien (Waffen, Bio/Chemie, Diskriminierung, Malware, EU-AI-Act-Verbote).",
  "Grenzfälle werden zur Prüfung an dich eskaliert — nicht still entschieden.","guardrails"),
 ("PHASE PLANEN","SCHRITT 6 VON 9","Der ZGPM-Plan entsteht",
  "Aus dem Verständnis bauen die Agenten einen vollständigen, methodischen Plan.",
  "Du verfolgst die Agenten und beurteilst den Plan: Gantt, RACI, Risk-Heatmap, Token-Live, Auslastung.",
  "PMO zerlegt in Phasen & Meilensteine; Worker füllen Rollen, Risiken, Aufwände; Reviewer prüft (max 3×).",
  "Bei Rot oder Konflikt entscheidest du.","plan"),
 ("PHASE PLANEN","SCHRITT 7 VON 9","Review & Edit — direkt am Bildschirm",
  "Du machst den Plan zu deinem Plan — Wortlaut, Termine, Werte.",
  "Text anklicken und ändern; Suffizienz bewusst bestätigen, bevor du freigibst.",
  "Es zeigt jede Änderung als Vorher/Nachher und versioniert sie (append-only).",
  "Erst deine Freigabe (Gate 2) macht aus der Version den gültigen Plan.","review"),
 ("PHASE BAUEN","SCHRITT 8 VON 9","Der Agenten-Harness wird gebaut & sichtbar",
  "Aus dem Plan wird ein lauffähiges Agententeam — visuell nachvollziehbar.",
  "Du betrachtest die Struktur, editierst Agenten (Mission/Skills) und ordnest per Kommando (Sequenz/Parallel/Skill/Agent).",
  "Es kompiliert Rollen zu Agenten, Aktivitäten zu Aufgaben, Risiken zu Quality-Gates; Reviewer flaggt Anti-Muster.",
  "Die HITL-Knoten ◆ zeigen, wo dein Sign-off verankert ist.","harness"),
 ("PHASE BAUEN","SCHRITT 9 VON 9","Export: signiertes Zip + Setup",
  "Du erhältst ein portables Paket, das ohne den Planner läuft.",
  "Freigeben (Gate 3), herunterladen, in Claude Cowork / Claude Code öffnen.",
  "Es schnürt CLAUDE.md, Skills, Agenten, Plan, Hooks & Handover in eine signierte Zip (checksums.txt).",
  "Die Prüfsumme bestätigt, dass dein Paket unverändert ist.","export"),
]
for phase,step,title,ziel,aktion,system,kontrolle,shot in STEPS:
    s=slide(); titlebar(s,f"{phase}   ·   {step}",title)
    tf=text(s,Inches(0.6),Inches(1.5),Inches(5.7),Inches(5.3))
    for lab,txt,col in [("ZIEL",ziel,NAVY),("DEINE AKTION",aktion,NAVY),
                        ("WAS DAS SYSTEM TUT",system,NAVY),("DEIN KONTROLLPUNKT",kontrolle,ORANGE)]:
        line_p(tf,lab,11.5,col,bold=True,after=1,before=6,first=(lab=="ZIEL"))
        line_p(tf,txt,12.5,INK,after=2)
    screen(s,(Inches(6.6),Inches(1.55),Inches(6.2),Inches(4.95)),f"{D}/{shot}.png")
    footer(s,pg())

# ============================================================ 15-16 DROPBOX
def dropbox_header(s,sub):
    titlebar(s,"PHASE PLANEN · SCHRITT 2A · CLOUD-QUELLE",f"Dropbox verbinden — Schritt für Schritt ({sub})")
def codebox(s,l,t,w,lines):
    h=Inches(0.3)*len(lines)+Inches(0.22)
    rect(s,l,t,w,h,NAVYD,rounded=True)
    tf=text(s,l+Inches(0.22),t+Inches(0.11),w-Inches(0.4),h-Inches(0.2))
    for i,ln in enumerate(lines): line_p(tf,ln,11,RGBColor(0xCF,0xD6,0xE6),first=(i==0),font=MONO,after=2)
    return t+h
def hint(s,l,t,w,txt):
    rect(s,l,t,w,Inches(1.0),VELLUM,line=GOLD,lw=Pt(1.25),rounded=True)
    tf=text(s,l+Inches(0.25),t+Inches(0.16),w-Inches(0.5),Inches(0.7))
    line_p(tf,"TRUST-HINWEIS",10,ORANGE,bold=True,first=True,after=3)
    line_p(tf,txt,11.5,INK,after=0)

s=slide(); dropbox_header(s,"1 / 2")
tf=text(s,Inches(0.6),Inches(1.5),Inches(7.2),Inches(5.2))
line_p(tf,"Ziel: einen Dropbox-Ordner als Quelle anbinden. Einmalige Einrichtung — danach genügt der Ordnerpfad.",13,STEEL,first=True,after=12)
line_p(tf,"App anlegen",14,NAVY,bold=True,after=6)
line_p(tf,"Auf dropbox.com/developers → App Console öffnen, „Create app“ klicken.",12.5,INK,bullet="1",after=6)
line_p(tf,"„Scoped access“ wählen · Zugriff „Full Dropbox“ (oder „App folder“) · App benennen.",12.5,INK,bullet="2",after=12)
line_p(tf,"Berechtigungen & Schlüssel",14,NAVY,bold=True,after=6)
line_p(tf,"Reiter „Permissions“: files.metadata.read + files.content.read aktivieren → „Submit“.",12.5,INK,bullet="3",after=6)
line_p(tf,"Reiter „Settings“: „App key“ + „App secret“ notieren (= DROPBOX_APP_KEY / _SECRET).",12.5,INK,bullet="4",after=6)
screen(s,(Inches(8.3),Inches(1.55),Inches(4.5),Inches(3.0)),f"{D}/interview.png")
hint(s,Inches(8.3),Inches(5.0),Inches(4.5),"Ohne Secrets bleibt der Connector ehrlich blockiert (kein Fake-Connect). SharePoint/OneDrive/Azure Blob sind bewusst gesperrt.")
footer(s,pg())

s=slide(); dropbox_header(s,"2 / 2")
tf=text(s,Inches(0.6),Inches(1.5),Inches(12.1),Inches(0.8))
line_p(tf,"Refresh-Token erzeugen (einmalig)",14,NAVY,bold=True,first=True,after=4)
line_p(tf,"URL im Browser öffnen (APP_KEY einsetzen), Zugriff erlauben, angezeigten Code kopieren:",12.5,INK,bullet="5",after=2)
y=codebox(s,Inches(0.95),Inches(2.7),Inches(11.5),["https://www.dropbox.com/oauth2/authorize?client_id=APP_KEY&response_type=code&token_access_type=offline"])
tf=text(s,Inches(0.6),y+Inches(0.14),Inches(12.1),Inches(0.5))
line_p(tf,"Code gegen Refresh-Token tauschen (Terminal) — Feld „refresh_token“ = DROPBOX_REFRESH_TOKEN:",12.5,INK,bullet="6",first=True,after=2)
y=codebox(s,Inches(0.95),y+Inches(0.78),Inches(11.5),["curl https://api.dropbox.com/oauth2/token \\","  -d code=CODE -d grant_type=authorization_code -u APP_KEY:APP_SECRET"])
tf=text(s,Inches(0.6),y+Inches(0.18),Inches(12.1),Inches(1.6))
line_p(tf,"In Betrieb nehmen",14,NAVY,bold=True,first=True,after=5)
line_p(tf,"DROPBOX_APP_KEY / _SECRET / _REFRESH_TOKEN in .env eintragen (nie echte Secrets ins Repo).",12.5,INK,bullet="7",after=6)
line_p(tf,"Im Planner (Schritt 2 · Cloud-Quelle): Dropbox zeigt „bereit“ → Ordnerpfad eingeben → „Ordner einlesen“. Inhalte ephemer; nur Name + Hash bleiben (Gate-1-Freeze).",12.5,INK,bullet="8",after=0)
footer(s,pg())

# ============================================================ 17 LOCKSCREEN
s=slide(); titlebar(s,"SICHERHEIT · HUMAN-IN-THE-LOOP","Sperrbildschirm — deine Sitzung bleibt geschützt")
tf=text(s,Inches(0.6),Inches(1.55),Inches(7.0),Inches(5.2))
line_p(tf,"Lässt du den Planner unbeaufsichtigt, sperrt er sich nach Inaktivität selbst — kein Bypass, der letzte Stand bleibt erhalten.",13,STEEL,first=True,after=12)
line_p(tf,"Was passiert",14,NAVY,bold=True,after=6)
line_p(tf,"Nach Idle-Timeout (Standard 15 min) legt sich ein Blur-Overlay über die App.",12.5,INK,bullet="1",after=6)
line_p(tf,"Während der Sperre setzt Aktivität den Timer nicht zurück — entsperrt wird nur bewusst.",12.5,INK,bullet="2",after=12)
line_p(tf,"So entsperrst du",14,NAVY,bold=True,after=6)
line_p(tf,"„Mit Single-Sign-On entsperren“ — danach läuft alles am letzten Stand weiter.",12.5,INK,bullet="3",after=6)
line_p(tf,"„Andere Person? Abmelden“ beendet die Sitzung für einen Nutzerwechsel.",12.5,INK,bullet="4",after=6)
screen(s,(Inches(8.0),Inches(1.7),Inches(4.8),Inches(3.3)),f"{D}/lockscreen.png")
hint(s,Inches(8.0),Inches(5.25),Inches(4.8),"Trust by design: Timeout konfigurierbar (NEXT_PUBLIC_LOCK_IDLE_SEC). SSO/MFA-Entsperrung folgt mit Entra-ID (Phase-2-Beta).")
footer(s,pg())

# ============================================================ 18 PROZESS / McK
s=slide(); titlebar(s,"PROZESS · NACH McKINSEY OPTIMIERT","Warum dieser Ablauf trägt")
tf=text(s,Inches(0.6),Inches(1.5),Inches(6.0),Inches(5.2))
for t,d in [("MECE-Phasen","Verstehen · Planen · Bauen — lückenlos, überschneidungsfrei."),
            ("Hypothesengeleitet","Das Interview testet Annahmen, statt nur zu sammeln."),
            ("Pyramid Principle","Jeder Meilenstein-Status nennt die Kernaussage zuerst."),
            ("Front-loaded Discovery","Projektart & Leitplanken vor der Planung geklärt."),
            ("Orchestrator statt Solo","PMO delegiert; Worker arbeiten parallel."),
            ("Evaluator-Loop mit Limit","Reviewer prüft max. 3×, dann entscheidest du."),
            ("Reversibilität by design","Jede Version bleibt erhalten — neu planen statt überschreiben.")]:
    line_p(tf,t,12.5,NAVY,bold=True,after=0,before=3,first=(t=="MECE-Phasen"))
    line_p(tf,d,12,INK,after=1)
c=rect(s,Inches(6.9),Inches(1.6),Inches(5.8),Inches(4.9),WHITE,line=LINE,rounded=True)
tf=text(s,Inches(7.2),Inches(1.9),Inches(5.2),Inches(4.4))
line_p(tf,"DREI GATES ◆ · DREI SCHLEIFEN ↺",12,ORANGE,bold=True,first=True,after=10)
for ph,st,g in [("VERSTEHEN","1 Beschreiben · 2 Interview ↺ · 3 Verständnis","◆ Gate 1"),
                ("PLANEN","4 Verwalten · 5 Leitplanken · 6 Plan ↺ · 7 Review","◆ Gate 2"),
                ("BAUEN","8 Harness ↺ · 9 Export","◆ Gate 3")]:
    line_p(tf,ph,12.5,NAVY,bold=True,after=1,before=4)
    line_p(tf,st,12,INK,after=1)
    line_p(tf,g,12,GREEN,bold=True,after=2)
line_p(tf,"Nichts geht ohne dich weiter. Schleifen: Interview, Reviewer (max 3×), Iteration.",12,STEEL,italic=True,after=0,before=6)
footer(s,pg())

# ============================================================ 19 BP ZGPM
s=slide(); titlebar(s,"BEST PRACTICES · ZGPM","Methodik nach Glasner et al. (PwC)")
tf=text(s,Inches(0.6),Inches(1.5),Inches(6.0),Inches(5.2))
line_p(tf,"DIE VIER BAUSTEINE",12,ORANGE,bold=True,first=True,after=8)
for t,d in [("Meilenstein","Zustand zu einem Termin (Verb im Perfekt) — „Konzept freigegeben“."),
            ("Aktivität","Arbeit vor dem Meilenstein, mit Aufwand & Verantwortlichen."),
            ("Ergebnispfad","Vertikaler Strang gleichartiger Ergebnisse (P/S/O)."),
            ("Phase","Zeitabschnitt, der Meilensteine bündelt (PH1·PH2·PH3).")]:
    line_p(tf,t,12.5,NAVY,bold=True,after=0,before=4); line_p(tf,d,12,INK,after=1)
tf=text(s,Inches(6.9),Inches(1.5),Inches(5.8),Inches(5.2))
line_p(tf,"PVM-CODES & HARTE REGELN",12,ORANGE,bold=True,first=True,after=8)
line_p(tf,"A führt aus · B beteiligt · E entscheidet · e entscheidet mit · F steuert Fortschritt · L leitet & steuert · I informiert · V verfügbar",12,INK,after=8)
line_p(tf,"Mindestens ein A pro Meilenstein/Aktivität.",12,INK,bullet="•",after=3)
line_p(tf,"Genau ein F oder L — nie mehr.",12,INK,bullet="•",after=3)
line_p(tf,"Ein „e“ steht nie allein, immer mit einem E.",12,INK,bullet="•",after=8)
line_p(tf,"RISIKO-AMPEL (propagiert nach oben)",12,ORANGE,bold=True,after=6)
line_p(tf,"grün — im Plan · gelb — achtgeben · rot — Lauf stoppt, du gibst frei",12,INK,after=0)
footer(s,pg())

# ============================================================ 20 BP HARNESS
s=slide(); titlebar(s,"BEST PRACTICES · AGENT-HARNESS (ANTHROPIC)","Muster & Anti-Muster")
tf=text(s,Inches(0.6),Inches(1.5),Inches(6.0),Inches(5.2))
line_p(tf,"PFLICHT-MUSTER",12,ORANGE,bold=True,first=True,after=7)
for t in ["Orchestrator-Worker — Lead zerlegt & delegiert mit Output-Schema.",
          "Evaluator-Optimizer — Reviewer prüft; max 3 Runden, dann HITL.",
          "Parallel-Tool-Calling — unabhängige Werkzeuge gleichzeitig.",
          "Filesystem-Artifact — große Ergebnisse als Datei, nur Referenz.",
          "Checkpoint & Resume — Zustand nach jedem Knoten sichern.",
          "Guardrails als eigener Prüf-Aufruf — nicht im Worker-Prompt.",
          "HITL an festen Punkten — Meilenstein, rotes Risiko, neue Skill, Budget."]:
    line_p(tf,t,12,INK,bullet="•",after=4)
tf=text(s,Inches(6.9),Inches(1.5),Inches(5.8),Inches(5.2))
line_p(tf,"ANTI-MUSTER (Reviewer flaggt hart)",12,ORANGE,bold=True,first=True,after=7)
for t in ["Vage Delegation ohne Ziel, Schema, Grenzen.","Über-Spawning — zu viele Subagenten.",
          "Routing im Prompt statt im Code.","Guardrail & Inhalt im selben Aufruf.",
          "Kein Checkpoint/Retry nach Fehlern.","Sequenziell statt parallel.",
          "Endlosschleifen ohne Stop-Bedingung.","Relative Pfade in zustandsbehafteten Agenten.",
          "Token-Budget ohne Prüfung."]:
    line_p(tf,t,12,INK,bullet="✕",after=4)
footer(s,pg())

# ============================================================ 21 GLOSSAR
s=slide(); titlebar(s,"REFERENZ","Glossar — ZGPM & Agentik")
tf=text(s,Inches(0.6),Inches(1.5),Inches(6.0),Inches(5.2))
line_p(tf,"ZGPM",12,ORANGE,bold=True,first=True,after=7)
for k,v in [("Meilenstein","Erreichter Zustand zu einem Termin."),("Aktivität","Arbeit vor dem Meilenstein."),
            ("Ergebnispfad","Strang gleichartiger Ergebnisse."),("MSP","Meilensteinplan."),
            ("PVM","Projektverantwortlichkeitsmatrix (RACI-Vorläufer)."),("PRL/MRL","Projekt-/Meilenstein-Risikoliste."),
            ("Ampel","Risiko-Status grün/gelb/rot, propagiert nach oben.")]:
    p=line_p(tf,k+" — ",12,NAVY,bold=True,after=3); r=p.add_run(); r.text=v; r.font.name=F; r.font.size=Pt(12); r.font.color.rgb=INK
tf=text(s,Inches(6.9),Inches(1.5),Inches(5.8),Inches(5.2))
line_p(tf,"AGENTIK",12,ORANGE,bold=True,first=True,after=7)
for k,v in [("Agent","Digitaler Mitarbeiter mit Rolle, Auftrag, Werkzeugen."),("Subagent","Agent mit isoliertem Kontext."),
            ("Skill","Wiederverwendbares Wissen/Werkzeug (SKILL.md)."),("Hook","Deterministische Regel (z. B. Stopp bei Rot)."),
            ("Orchestrator","Lead-Agent, der zerlegt & delegiert (PMO)."),("HITL","Human-in-the-Loop — du als Freigabe-Instanz."),
            ("Harness","Portables Paket: CLAUDE.md, Skills, Agenten, Plan."),("Cowork/Claude Code","Umgebung, in der das Harness läuft.")]:
    p=line_p(tf,k+" — ",12,NAVY,bold=True,after=3); r=p.add_run(); r.text=v; r.font.name=F; r.font.size=Pt(12); r.font.color.rgb=INK
footer(s,pg())

# ============================================================ 22 SCHNELLREFERENZ
s=slide(); titlebar(s,"REFERENZ","Der ganze Weg auf einer Seite")
tf=text(s,Inches(0.6),Inches(1.5),Inches(7.4),Inches(5.2))
line_p(tf,"DIE NEUN SCHRITTE",12,ORANGE,bold=True,first=True,after=7)
for i,t in enumerate(["Projekt in eigenen Worten beschreiben.","Schärfungs-Interview beantworten (McK).",
   "Verständnis & Agentenstruktur freigeben.","Projekte verwalten — speichern, kopieren, löschen.",
   "Leitplanken kennen — was geht, was nicht.","ZGPM-Plan lesen: Gantt, RACI, Risiko, Token.",
   "Am Bildschirm anpassen — Text & Werte.","Harness mitgestalten — Artefakte & Kommandos.",
   "Als Zip exportieren → Claude Cowork / Code."],1):
    line_p(tf,t,12.5,INK,bullet=str(i),after=4)
c=rect(s,Inches(8.3),Inches(1.6),Inches(4.4),Inches(4.9),NAVY,rounded=True)
tf=text(s,Inches(8.6),Inches(1.9),Inches(3.85),Inches(4.4))
line_p(tf,"DEINE 3 KONTROLLPUNKTE",11.5,GOLD,bold=True,first=True,after=6)
for t in ["Meilenstein-Freigabe","Rote Risiko-Ampel","Harness-Freigabe (Gate 3)"]:
    line_p(tf,t,12,WHITE,bullet="◆",after=3)
line_p(tf,"WAS NIE GEHT",11.5,GOLD,bold=True,after=6,before=8)
for t in ["Waffen & Dual-Use","Bio/Chemie/Nuklear","Diskriminierung","Malware & Exploits","EU-AI-Act-Verbote"]:
    line_p(tf,t,12,WHITE,bullet="✕",after=3)
footer(s,pg())

# ============================================================ 23 ABSCHLUSS
s=slide(); rect(s,0,0,SW,SH,NAVY); rect(s,0,Inches(4.4),SW,Inches(0.08),GOLD)
tf=text(s,Inches(0.9),Inches(2.0),Inches(11.5),Inches(2.2))
line_p(tf,"Du startest mit Worten.",34,WHITE,bold=True,first=True,after=4)
line_p(tf,"Du endest mit einem lauffähigen Agententeam.",34,GOLD,bold=True,after=0)
tf=text(s,Inches(0.9),Inches(4.8),Inches(11.5),Inches(1.6))
line_p(tf,"1. Beschreiben — öffne den Planner und beschreibe dein Projekt.",14,RGBColor(0xC8,0xD2,0xEC),first=True,after=4)
line_p(tf,"2. Freigeben — schärfe das Verständnis und gib es frei.",14,RGBColor(0xC8,0xD2,0xEC),after=4)
line_p(tf,"3. Bauen — lass den Harness bauen und exportiere ihn.",14,RGBColor(0xC8,0xD2,0xEC),after=0)

REPO="/home/user/agent_operating_model_planner/user_guides/AEGIRA_Planner_User_Guide.pptx"
prs.save(REPO); prs.save("/tmp/AEGIRA_Planner_User_Guide_CLEAN.pptx")
print("Folien:", len(prs.slides._sldIdLst), "-> gespeichert")
